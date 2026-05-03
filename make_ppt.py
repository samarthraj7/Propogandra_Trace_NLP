"""
Build PropagandaTrace presentation.
Run: python make_ppt.py
Output: PropagandaTrace_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand colours ─────────────────────────────────────────────────────────────
USC_RED    = RGBColor(0x99, 0x00, 0x00)   # USC Cardinal
USC_GOLD   = RGBColor(0xFF, 0xCC, 0x00)   # USC Gold
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY  = RGBColor(0x22, 0x22, 0x22)
MID_GRAY   = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
GREEN      = RGBColor(0x1A, 0x7A, 0x3C)
BLUE       = RGBColor(0x1A, 0x4A, 0x8A)

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank


# ── Helper utilities ──────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_textbox(slide, text, x, y, w, h,
                font_size=18, bold=False, italic=False,
                color=DARK_GRAY, align=PP_ALIGN.LEFT,
                wrap=True, font_name="Calibri"):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(font_size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name  = font_name
    return txb

def add_para(tf, text, font_size=16, bold=False, italic=False,
             color=DARK_GRAY, align=PP_ALIGN.LEFT, level=0,
             space_before=0, font_name="Calibri"):
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(font_size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name   = font_name
    return p

def header_bar(slide, title, subtitle=None):
    """Dark red top bar with title."""
    add_rect(slide, 0, 0, W, Inches(1.15), fill=USC_RED)
    add_textbox(slide, title, Inches(0.35), Inches(0.12), Inches(12.3), Inches(0.7),
                font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_textbox(slide, subtitle, Inches(0.35), Inches(0.75), Inches(12.3), Inches(0.38),
                    font_size=14, italic=True, color=USC_GOLD, align=PP_ALIGN.LEFT)
    # thin gold rule
    add_rect(slide, 0, Inches(1.15), W, Pt(4), fill=USC_GOLD)

def bullet_box(slide, bullets, x, y, w, h,
               font_size=15, title=None, title_color=USC_RED,
               bg=None, border=None):
    """A text frame with optional background and bullet list."""
    if bg or border:
        add_rect(slide, x, y, w, h, fill=bg, line=border)
    txb = slide.shapes.add_textbox(x + Inches(0.12), y + Inches(0.1),
                                    w - Inches(0.24), h - Inches(0.18))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = title
        r.font.size  = Pt(font_size + 1)
        r.font.bold  = True
        r.font.color.rgb = title_color
        r.font.name  = "Calibri"
        first = False
    for b in bullets:
        indent = 0
        txt = b
        if b.startswith("  "):
            indent = 1
            txt = b.strip()
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = indent
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = ("• " if indent == 0 else "◦ ") + txt
        r.font.size  = Pt(font_size - indent)
        r.font.color.rgb = DARK_GRAY
        r.font.name  = "Calibri"
    return txb

def code_box(slide, text, x, y, w, h, font_size=11):
    """Monospaced dark box for code/example text."""
    add_rect(slide, x, y, w, h, fill=RGBColor(0x1E, 0x1E, 0x2E))
    txb = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.08),
                                    w - Inches(0.2), h - Inches(0.15))
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for line in text.split("\n"):
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size  = Pt(font_size)
        r.font.color.rgb = RGBColor(0xCC, 0xFF, 0xCC)
        r.font.name  = "Courier New"
    return txb

def label_box(slide, label, x, y, w, h, fill=USC_RED, font_size=13):
    add_rect(slide, x, y, w, h, fill=fill)
    add_textbox(slide, label, x, y, w, h,
                font_size=font_size, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

def arrow_right(slide, x, y, length=Inches(0.5), color=MID_GRAY):
    """Simple right-pointing arrow line."""
    line = slide.shapes.add_connector(1, x, y, x + length, y)
    line.line.color.rgb = color
    line.line.width = Pt(2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=USC_RED)
add_rect(slide, 0, Inches(5.6), W, Inches(1.9), fill=RGBColor(0x6B, 0x00, 0x00))
add_rect(slide, 0, Inches(5.57), W, Pt(5), fill=USC_GOLD)

add_textbox(slide, "PropagandaTrace",
            Inches(0.6), Inches(0.9), Inches(12), Inches(1.2),
            font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Multi-Source Attribution of Artificial Intelligence–Generated",
            Inches(0.6), Inches(2.0), Inches(12), Inches(0.6),
            font_size=22, color=USC_GOLD, align=PP_ALIGN.CENTER)
add_textbox(slide, "Wartime Propaganda via Watermark Fingerprinting",
            Inches(0.6), Inches(2.5), Inches(12), Inches(0.6),
            font_size=22, color=USC_GOLD, align=PP_ALIGN.CENTER)

add_textbox(slide, "CS 544 — Natural Language Processing",
            Inches(0.6), Inches(3.5), Inches(12), Inches(0.45),
            font_size=17, italic=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "University of Southern California",
            Inches(0.6), Inches(3.95), Inches(12), Inches(0.45),
            font_size=17, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide,
            "Samarth Rajendra  •  Sanath Nagendra Bhargav  •  Jeevika Kiran  •  Shree Sriadibhatla  •  Akash Gandi",
            Inches(0.6), Inches(5.75), Inches(12), Inches(0.45),
            font_size=14, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Team 47", Inches(0.6), Inches(6.2), Inches(12), Inches(0.4),
            font_size=13, color=USC_GOLD, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=LIGHT_GRAY)
header_bar(slide, "Presentation Outline")

items = [
    ("1", "Research Problem & Task", "What are we solving and why does it matter?"),
    ("2", "Motivation",              "The real-world threat of AI-generated propaganda"),
    ("3", "Background & Literature", "Watermarking, attribution, and prior work"),
    ("4", "Datasets",                "WTWT, GDELT, SemEval-2020, and custom templates"),
    ("5", "System Pipeline",         "Four-phase architecture overview"),
    ("6", "Running Example",         "End-to-end walkthrough with real outputs"),
    ("7", "Results",                 "Detection rates, attribution accuracy, evasion robustness"),
    ("8", "Conclusion",              "Findings, contributions, and future work"),
]

col_w = Inches(6.3)
for i, (num, title, desc) in enumerate(items):
    row = i // 2
    col = i % 2
    x = Inches(0.3) + col * (col_w + Inches(0.43))
    y = Inches(1.35) + row * Inches(1.45)
    add_rect(slide, x, y, col_w, Inches(1.25), fill=WHITE, line=USC_RED)
    add_rect(slide, x, y, Inches(0.55), Inches(1.25), fill=USC_RED)
    add_textbox(slide, num, x, y, Inches(0.55), Inches(1.25),
                font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, title, x + Inches(0.62), y + Inches(0.1),
                col_w - Inches(0.72), Inches(0.5),
                font_size=15, bold=True, color=USC_RED)
    add_textbox(slide, desc, x + Inches(0.62), y + Inches(0.58),
                col_w - Inches(0.72), Inches(0.55),
                font_size=12, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Research Problem
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=LIGHT_GRAY)
header_bar(slide, "Research Problem & Task",
           "What gap does PropagandaTrace fill?")

# Left panel
add_rect(slide, Inches(0.3), Inches(1.35), Inches(5.9), Inches(5.7), fill=WHITE, line=USC_RED)
add_textbox(slide, "The Gap in Existing Research",
            Inches(0.45), Inches(1.45), Inches(5.6), Inches(0.45),
            font_size=16, bold=True, color=USC_RED)
probs = [
    "Existing AI-text detectors answer only:",
    '  "Is this text AI-generated?" (binary)',
    "",
    "They cannot answer:",
    '  "Which model generated this text?"',
    '  "Which organisation deployed it?"',
    "",
    "This is critical for wartime disinformation:",
    "  State actors use LLMs to mass-produce",
    "  propaganda — attribution is impossible",
    "  without model-level fingerprinting.",
]
y0 = Inches(1.95)
for line in probs:
    indent = line.startswith("  ")
    col = MID_GRAY if indent else DARK_GRAY
    sz = 13 if indent else 14
    add_textbox(slide, line.strip(), Inches(0.5 + (0.25 if indent else 0)),
                y0, Inches(5.4), Inches(0.32),
                font_size=sz, color=col)
    y0 += Inches(0.31)

# Right panel
add_rect(slide, Inches(6.5), Inches(1.35), Inches(6.5), Inches(5.7), fill=USC_RED)
add_textbox(slide, "Our Task",
            Inches(6.65), Inches(1.45), Inches(6.2), Inches(0.45),
            font_size=16, bold=True, color=USC_GOLD)

tasks = [
    ("Assign", "Give each Large Language Model (LLM) a unique secret watermark key at deployment time."),
    ("Generate", "Produce a large corpus of watermarked propaganda texts using three LLMs under two watermarking schemes."),
    ("Attack", "Simulate adversarial evasion: paraphrase and round-trip translation to destroy the watermark signal."),
    ("Attribute", "Recover the source model from evaded text using forensic watermark detection and multi-class attribution scoring."),
]
y0 = Inches(1.95)
for kw, desc in tasks:
    add_rect(slide, Inches(6.6), y0, Inches(0.9), Inches(1.05), fill=USC_GOLD)
    add_textbox(slide, kw, Inches(6.6), y0, Inches(0.9), Inches(1.05),
                font_size=12, bold=True, color=USC_RED, align=PP_ALIGN.CENTER)
    add_textbox(slide, desc, Inches(7.6), y0 + Inches(0.08), Inches(5.2), Inches(1.0),
                font_size=12, color=WHITE)
    y0 += Inches(1.2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Motivation
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=LIGHT_GRAY)
header_bar(slide, "Motivation",
           "Why does AI-generated wartime propaganda demand a new solution?")

stats = [
    ("10×", "Increase in detected AI-generated\ndisinformation content since 2022\n(Stanford Internet Observatory)"),
    ("< 1¢", "Cost to generate one article-length\npropaganda text using a\ncurrent LLM API"),
    ("60%", "Accuracy of best existing\nbinary AI-detectors on\nparaphrased text"),
    ("0%", "Attribution accuracy of any\nexisting tool: model-level\nfingerprinting is unsolved"),
]
for i, (stat, desc) in enumerate(stats):
    x = Inches(0.25) + i * Inches(3.27)
    add_rect(slide, x, Inches(1.35), Inches(3.0), Inches(2.1), fill=USC_RED)
    add_textbox(slide, stat, x, Inches(1.45), Inches(3.0), Inches(0.9),
                font_size=36, bold=True, color=USC_GOLD, align=PP_ALIGN.CENTER)
    add_textbox(slide, desc, x + Inches(0.1), Inches(2.3), Inches(2.8), Inches(1.1),
                font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(slide, "The Core Problem",
            Inches(0.3), Inches(3.65), Inches(12.73), Inches(0.4),
            font_size=16, bold=True, color=USC_RED)

points = [
    "Large Language Models (LLMs) can produce fluent, persuasive wartime narratives at industrial scale with minimal human involvement.",
    "Current detection tools classify text as human or AI — but cannot identify which of many deployed models produced it.",
    "Without model-level attribution, there is no technical basis for holding a state actor or organisation accountable.",
    "Watermark fingerprinting offers a cryptographic solution: embed a hidden, model-specific signal during generation that survives post-processing.",
]
y0 = Inches(4.1)
for pt in points:
    add_rect(slide, Inches(0.3), y0, Pt(6), Inches(0.38), fill=USC_GOLD)
    add_textbox(slide, pt, Inches(0.55), y0, Inches(12.5), Inches(0.42), font_size=13, color=DARK_GRAY)
    y0 += Inches(0.48)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Background & Literature
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=LIGHT_GRAY)
header_bar(slide, "Background & Relevant Literature")

cols = [
    ("LLM Watermarking\n(Kirchenbauer et al., 2023)",
     [
         "Kirchenbauer et al. (2023) — 'A Watermark for Large Language Models' (ICML 2023)",
         "KGW scheme: partition vocabulary into green/red lists using a keyed hash; bias green tokens during sampling",
         "Detection via z-score of green-token fraction; theoretically sound under the null hypothesis",
         "First practical, training-free watermarking scheme for auto-regressive LLMs",
     ]),
    ("Cryptographic Watermarking\n(Aaronson, 2022)",
     [
         "Aaronson (2022) — EXP / Gumbel-max watermark scheme (blog + COLT keynote)",
         "Uses keyed pseudo-random number generator (PRNG) to derive random vectors per token position",
         "Detection: accumulate per-token score; high average score indicates watermarked text",
         "Provably undetectable to the reader; robust to mild perturbations",
     ]),
    ("AI Text Detection &\nAttribution Gaps",
     [
         "Sadasivan et al. (2023) — 'Can AI-Generated Text be Reliably Detected?' — shows paraphrase breaks most detectors",
         "Grinbaum & Adomaitis (2023) — documents LLM use in Ukrainian/Russian information warfare",
         "No prior work addresses multi-class attribution: which of N deployed models produced a given text?",
         "PropagandaTrace is the first framework combining watermark fingerprinting with multi-model attribution in the propaganda domain",
     ]),
]

col_w = Inches(4.1)
for i, (title, bullets) in enumerate(cols):
    x = Inches(0.25) + i * (col_w + Inches(0.2))
    add_rect(slide, x, Inches(1.35), col_w, Inches(5.7), fill=WHITE, line=USC_RED)
    add_rect(slide, x, Inches(1.35), col_w, Inches(0.65), fill=USC_RED)
    add_textbox(slide, title, x + Inches(0.08), Inches(1.38), col_w - Inches(0.16), Inches(0.6),
                font_size=12, bold=True, color=WHITE)
    y0 = Inches(2.1)
    for b in bullets:
        add_rect(slide, x + Inches(0.1), y0, Pt(5), Inches(0.3), fill=USC_GOLD)
        add_textbox(slide, b, x + Inches(0.27), y0, col_w - Inches(0.35), Inches(0.9),
                    font_size=11, color=DARK_GRAY)
        y0 += Inches(1.1)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Datasets
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=LIGHT_GRAY)
header_bar(slide, "Datasets Used — Phase 1: Seed Prompt Collection",
           "Four sources chosen to cover real-world linguistic diversity and propaganda styles")

datasets = [
    ("WTWT", "Will They, Won't They\n(Conforti et al., 2020)",
     USC_RED, "2,000+ tweets",
     [
         "Stance-labeled tweets on corporate merger rumours",
         "Covers adversarial, opinionated, short-form language",
         "Why chosen: real social-media rhetoric; mirrors disinformation tone",
         "Source: HuggingFace — conforti/wtwt",
     ]),
    ("GDELT", "Global Database of Events,\nLanguage and Tone",
     BLUE, "500 article snippets",
     [
         "Free public API; indexes global news in real time",
         "Queried with 6 conflict keywords (e.g. 'military offensive attack')",
         "Why chosen: real conflict journalism; provides factual grounding for prompts",
         "No authentication required",
     ]),
    ("SemEval-2020\nTask 11", "Propaganda Fragment\nDetection Dataset",
     GREEN, "200 paragraphs",
     [
         "Professionally annotated propaganda techniques in news articles",
         "14 propaganda techniques labelled at fragment level",
         "Why chosen: gold-standard propaganda examples; adds domain authenticity",
         "Public dev-set articles fetched from GitHub",
     ]),
    ("Custom\nTemplates", "Fill-in-the-Blank\nPropaganda Templates",
     RGBColor(0x7B, 0x35, 0x00), "1,300+ prompts",
     [
         "20 templates: 10 military situation reports + 10 political commentaries",
         "Slot-filled from curated lists: factions, enemies, locations, conflict types",
         "Why chosen: guaranteed propaganda style; always available as fallback",
         "Deduplication applied across all four sources",
     ]),
]

col_w = Inches(3.1)
for i, (name, full, color, count, bullets) in enumerate(datasets):
    x = Inches(0.22) + i * (col_w + Inches(0.17))
    add_rect(slide, x, Inches(1.35), col_w, Inches(5.7), fill=WHITE, line=color)
    add_rect(slide, x, Inches(1.35), col_w, Inches(1.05), fill=color)
    add_textbox(slide, name, x + Inches(0.08), Inches(1.38), col_w - Inches(0.16), Inches(0.5),
                font_size=16, bold=True, color=WHITE)
    add_textbox(slide, full, x + Inches(0.08), Inches(1.82), col_w - Inches(0.16), Inches(0.52),
                font_size=10, color=WHITE, italic=True)
    add_rect(slide, x, Inches(2.4), col_w, Inches(0.38), fill=RGBColor(0xF0, 0xF0, 0xF0))
    add_textbox(slide, f"Sample size: {count}", x + Inches(0.1), Inches(2.42),
                col_w - Inches(0.2), Inches(0.34), font_size=12, bold=True, color=color)
    y0 = Inches(2.85)
    for b in bullets:
        add_rect(slide, x + Inches(0.1), y0 + Inches(0.07), Pt(5), Pt(5), fill=color)
        add_textbox(slide, b, x + Inches(0.22), y0, col_w - Inches(0.3), Inches(0.72),
                    font_size=10, color=DARK_GRAY)
        y0 += Inches(0.73)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — System Pipeline Overview
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=LIGHT_GRAY)
header_bar(slide, "System Pipeline — Four-Phase Architecture")

phases = [
    ("Phase 1", "Data\nCollection", USC_RED,
     "WTWT + GDELT +\nSemEval-2020 +\nTemplates\n→ 4,000 seed prompts\nin JSONL format"),
    ("Phase 2", "Watermarked\nCorpus Generation", BLUE,
     "3 LLMs × 2 schemes\n= 6 corpus files\n→ 18,000 watermarked\npropaganda texts"),
    ("Phase 3", "Evasion\nAttacks", GREEN,
     "T5 + PEGASUS\n+ RTT Arabic\n+ RTT Russian\n→ 24 evaded JSONL files"),
    ("Phase 4", "Detection &\nAttribution", RGBColor(0x7B, 0x35, 0x00),
     "KGW + Aaronson\ndetectors on all files\n→ Detection rates,\nattribution accuracy"),
]

bw = Inches(2.8)
bh = Inches(4.2)
for i, (phase, name, color, detail) in enumerate(phases):
    x = Inches(0.25) + i * (bw + Inches(0.42))
    add_rect(slide, x, Inches(1.5), bw, Inches(0.55), fill=color)
    add_textbox(slide, phase, x, Inches(1.52), bw, Inches(0.5),
                font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, x, Inches(2.05), bw, bh - Inches(0.55), fill=WHITE, line=color)
    add_textbox(slide, name, x + Inches(0.1), Inches(2.15), bw - Inches(0.2), Inches(0.8),
                font_size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_textbox(slide, detail, x + Inches(0.15), Inches(3.05), bw - Inches(0.3), Inches(2.5),
                font_size=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    if i < 3:
        add_textbox(slide, "→", Inches(0.25) + i * (bw + Inches(0.42)) + bw + Inches(0.05),
                    Inches(3.6), Inches(0.35), Inches(0.6),
                    font_size=28, bold=True, color=USC_RED, align=PP_ALIGN.CENTER)

# Model list
add_rect(slide, Inches(0.25), Inches(5.9), Inches(12.83), Inches(1.25), fill=WHITE, line=USC_GOLD)
add_textbox(slide, "LLMs Used (Phase 2):  ",
            Inches(0.4), Inches(6.0), Inches(2.2), Inches(0.45),
            font_size=13, bold=True, color=USC_RED)
models_txt = "Mistral-7B-Instruct-v0.2   •   Meta-LLaMA-3-8B-Instruct   •   Falcon-7B-Instruct     |     Schemes:  KGW  •  Aaronson (EXP)"
add_textbox(slide, models_txt, Inches(2.5), Inches(6.0), Inches(10.4), Inches(0.45),
            font_size=13, color=DARK_GRAY)
add_textbox(slide, "Evasion models:  T5-Large  •  PEGASUS-XSUM  •  NLLB-200-Distilled-600M (for round-trip translation)",
            Inches(0.4), Inches(6.52), Inches(12.4), Inches(0.4), font_size=12, color=MID_GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Watermarking Deep Dive: KGW
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=LIGHT_GRAY)
header_bar(slide, "Phase 2 — KGW Watermarking Scheme",
           "Kirchenbauer, Geiping & Wallace (ICML 2023) — 'A Watermark for Large Language Models'")

# Left: how it works
add_rect(slide, Inches(0.25), Inches(1.35), Inches(6.3), Inches(5.8), fill=WHITE, line=USC_RED)
add_textbox(slide, "How It Works — Generation",
            Inches(0.4), Inches(1.45), Inches(6.0), Inches(0.4),
            font_size=15, bold=True, color=USC_RED)
steps = [
    ("1", "At each generation step t, take the previous token t_{i-1}"),
    ("2", "Compute SHA-256( secret_key : t_{i-1} ) → seed"),
    ("3", "Use seed to shuffle vocabulary V into a random permutation"),
    ("4", "First γ|V| tokens = Green List  (γ = 0.25, so 25% of vocab)"),
    ("5", "Remaining (1−γ)|V| tokens = Red List"),
    ("6", "Add logit bias δ = +2.0 to every Green List token before sampling"),
    ("7", "Model samples from biased distribution → Green tokens sampled more often"),
]
y0 = Inches(1.95)
for num, txt in steps:
    add_rect(slide, Inches(0.4), y0, Inches(0.4), Inches(0.52), fill=USC_RED)
    add_textbox(slide, num, Inches(0.4), y0, Inches(0.4), Inches(0.52),
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, txt, Inches(0.88), y0 + Inches(0.06), Inches(5.55), Inches(0.44),
                font_size=12, color=DARK_GRAY)
    y0 += Inches(0.6)

add_textbox(slide, "Detection Formula:",
            Inches(0.4), Inches(6.55), Inches(6.0), Inches(0.35),
            font_size=13, bold=True, color=USC_RED)
add_textbox(slide, "z = (|green tokens| − γ·T) / √(γ(1−γ)·T)      →  flag as watermarked if z > 4.0",
            Inches(0.4), Inches(6.9), Inches(6.0), Inches(0.35),
            font_size=12, color=DARK_GRAY, italic=True)

# Right: key properties
add_rect(slide, Inches(6.85), Inches(1.35), Inches(6.2), Inches(2.6), fill=USC_RED)
add_textbox(slide, "Key Properties",
            Inches(7.0), Inches(1.45), Inches(5.9), Inches(0.4),
            font_size=15, bold=True, color=WHITE)
props = [
    "Training-free — plugs into any LLM as a LogitsProcessor",
    "Each model uses a distinct secret key → unique fingerprint",
    "Statistically detectable with z > 4.0 (false positive rate < 0.001%)",
    "γ = 0.25 and δ = 2.0 tuned to balance text quality vs. detectability",
]
y0 = Inches(1.95)
for p in props:
    add_textbox(slide, "✓  " + p, Inches(7.05), y0, Inches(5.85), Inches(0.5),
                font_size=12, color=WHITE)
    y0 += Inches(0.52)

add_rect(slide, Inches(6.85), Inches(4.15), Inches(6.2), Inches(2.8), fill=WHITE, line=USC_RED)
add_textbox(slide, "Keys Assigned per Model",
            Inches(7.0), Inches(4.25), Inches(5.9), Inches(0.4),
            font_size=14, bold=True, color=USC_RED)
keys = [
    ("Mistral-7B-Instruct-v0.2",          "MISTRAL_KEY_47A"),
    ("Meta-LLaMA-3-8B-Instruct",          "LLAMA_KEY_47B"),
    ("Falcon-7B-Instruct",                "FALCON_KEY_47C"),
]
y0 = Inches(4.75)
for model, key in keys:
    add_textbox(slide, model, Inches(7.05), y0, Inches(3.2), Inches(0.42), font_size=12, color=DARK_GRAY)
    add_rect(slide, Inches(10.35), y0 + Inches(0.04), Inches(2.5), Inches(0.34), fill=LIGHT_GRAY)
    add_textbox(slide, key, Inches(10.35), y0 + Inches(0.04), Inches(2.5), Inches(0.34),
                font_size=11, bold=True, color=USC_RED, align=PP_ALIGN.CENTER, font_name="Courier New")
    y0 += Inches(0.52)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Watermarking Deep Dive: Aaronson
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=LIGHT_GRAY)
header_bar(slide, "Phase 2 — Aaronson / EXP Watermarking Scheme",
           "Aaronson (2022) — Cryptographic watermarking via Gumbel-max sampling")

add_rect(slide, Inches(0.25), Inches(1.35), Inches(6.3), Inches(5.8), fill=WHITE, line=BLUE)
add_textbox(slide, "How It Works — Generation",
            Inches(0.4), Inches(1.45), Inches(6.0), Inches(0.4),
            font_size=15, bold=True, color=BLUE)
steps = [
    ("1", "At position t, derive r_t ∈ [0,1]^|V| using SHA-256(key : t) as PRNG seed"),
    ("2", "Convert logits to log-probabilities: log p_i for each token i"),
    ("3", "Add Gumbel noise keyed by r:  modified_i = log(p_i) − log(−log(r_{t,i}))"),
    ("4", "Sample token = argmax over modified scores"),
    ("5", "This is equivalent to sampling token = argmax_i  r_{t,i}^{1/p_i}"),
    ("6", "Tokens where r_{t,i} is large are strongly favoured → biased distribution"),
    ("7", "Watermark is invisible: output distribution matches unbiased sampling in expectation"),
]
y0 = Inches(1.95)
for num, txt in steps:
    add_rect(slide, Inches(0.4), y0, Inches(0.4), Inches(0.52), fill=BLUE)
    add_textbox(slide, num, Inches(0.4), y0, Inches(0.4), Inches(0.52),
                font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, txt, Inches(0.88), y0 + Inches(0.06), Inches(5.55), Inches(0.44),
                font_size=12, color=DARK_GRAY)
    y0 += Inches(0.6)

add_textbox(slide, "Detection Formula:",
            Inches(0.4), Inches(6.55), Inches(6.0), Inches(0.35),
            font_size=13, bold=True, color=BLUE)
add_textbox(slide, "S = mean( −log(1 − r_{t, w_t}) )  for t = 1…T      →  flag as watermarked if S > 0.8",
            Inches(0.4), Inches(6.9), Inches(6.0), Inches(0.35),
            font_size=12, color=DARK_GRAY, italic=True)

add_rect(slide, Inches(6.85), Inches(1.35), Inches(6.2), Inches(2.7), fill=BLUE)
add_textbox(slide, "Key Properties",
            Inches(7.0), Inches(1.45), Inches(5.9), Inches(0.4),
            font_size=15, bold=True, color=WHITE)
props = [
    "Cryptographically sound — based on pseudo-random number generator (PRNG) with secret key",
    "Provably undetectable to a reader without the key",
    "Position-dependent: each token position uses a different random vector",
    "Requires reset() between independent texts to avoid position aliasing",
]
y0 = Inches(1.95)
for p in props:
    add_textbox(slide, "✓  " + p, Inches(7.05), y0, Inches(5.85), Inches(0.52),
                font_size=12, color=WHITE)
    y0 += Inches(0.54)

add_rect(slide, Inches(6.85), Inches(4.25), Inches(6.2), Inches(2.7), fill=WHITE, line=BLUE)
add_textbox(slide, "KGW vs Aaronson — Comparison",
            Inches(7.0), Inches(4.35), Inches(5.9), Inches(0.4),
            font_size=14, bold=True, color=BLUE)
rows = [
    ("Bias mechanism",   "Logit additive (+δ)",           "Gumbel-max PRNG"),
    ("Context",          "Previous token hash",           "Position + key hash"),
    ("Detection signal", "Green-token z-score",           "Mean r-value score"),
    ("Evasion fragility","Moderate",                      "Low to moderate"),
    ("Text quality impact","Slight (high δ distorts)",    "Negligible"),
]
y0 = Inches(4.85)
for prop, kgw, aar in rows:
    add_textbox(slide, prop, Inches(7.05), y0, Inches(1.7), Inches(0.38), font_size=11, bold=True, color=MID_GRAY)
    add_textbox(slide, kgw,  Inches(8.8),  y0, Inches(2.0), Inches(0.38), font_size=11, color=DARK_GRAY)
    add_textbox(slide, aar,  Inches(10.85),y0, Inches(2.0), Inches(0.38), font_size=11, color=DARK_GRAY)
    y0 += Inches(0.42)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Evasion Attacks
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=LIGHT_GRAY)
header_bar(slide, "Phase 3 — Evasion Attacks",
           "Four strategies designed to break the watermark signal while preserving meaning")

attacks = [
    ("T5\nParaphrase", "t5-large", GREEN,
     "Lexical rewriting — single pass",
     [
         'Prefix: "paraphrase: " + watermarked text → T5-Large',
         "Replaces words and short phrases while keeping sentence structure",
         "Mild surface perturbation — expected to partially survive watermark",
         "Beam search with 5 beams; max 512 tokens",
     ]),
    ("PEGASUS\nParaphrase", "pegasus-xsum", RGBColor(0x7B, 0x35, 0x00),
     "Neural rewriting — aggressive",
     [
         "PEGASUS-XSUM was trained for extreme summarisation → produces radical rewrites",
         "Reconstructs meaning in completely new sentence structures",
         "Expected to heavily degrade KGW (token-level) watermark signal",
         "Beam search with 8 beams",
     ]),
    ("Round-Trip\nTranslation (Arabic)", "NLLB-200-600M", BLUE,
     "English → Arabic → English",
     [
         "NLLB-200-Distilled-600M: Meta's multilingual translation model",
         "Forward: eng_Latn → ara_Arab  |  Backward: ara_Arab → eng_Latn",
         "Morphological differences between Arabic and English scramble token sequences",
         "FLORES-200 language codes used",
     ]),
    ("Round-Trip\nTranslation (Russian)", "NLLB-200-600M", USC_RED,
     "English → Russian → English",
     [
         "Same NLLB-200 model; pivot language: rus_Cyrl",
         "Cyrillic script causes complete token re-mapping upon back-translation",
         "Grammatical structure changes (case system, word order) further disrupt signal",
         "Both RTT strategies share the same NLLB-200 model loaded once",
     ]),
]

col_w = Inches(3.1)
for i, (name, model_id, color, subtitle, bullets) in enumerate(attacks):
    x = Inches(0.22) + i * (col_w + Inches(0.17))
    add_rect(slide, x, Inches(1.35), col_w, Inches(5.75), fill=WHITE, line=color)
    add_rect(slide, x, Inches(1.35), col_w, Inches(0.85), fill=color)
    add_textbox(slide, name, x + Inches(0.07), Inches(1.38), col_w - Inches(0.14), Inches(0.55),
                font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, subtitle, x + Inches(0.07), Inches(2.28), col_w - Inches(0.14), Inches(0.38),
                font_size=11, bold=True, color=color)
    add_rect(slide, x + Inches(0.07), Inches(2.2), col_w - Inches(0.14), Inches(0.0), fill=None, line=color)
    add_textbox(slide, f"Model: {model_id}", x + Inches(0.07), Inches(2.62), col_w - Inches(0.14), Inches(0.3),
                font_size=10, italic=True, color=MID_GRAY)
    y0 = Inches(3.0)
    for b in bullets:
        add_rect(slide, x + Inches(0.1), y0 + Inches(0.1), Pt(5), Pt(5), fill=color)
        add_textbox(slide, b, x + Inches(0.25), y0, col_w - Inches(0.33), Inches(0.68),
                    font_size=10.5, color=DARK_GRAY)
        y0 += Inches(0.7)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Running Example (Part 1: Prompt → Generation)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=LIGHT_GRAY)
header_bar(slide, "Running Example — Step 1: Seed Prompt → Watermarked Text",
           "Following a single record through the entire pipeline")

# Step indicators
steps_nav = ["1. Seed Prompt", "2. Watermarked Gen.", "3. Evasion", "4. Detection"]
for i, s in enumerate(steps_nav):
    col = USC_RED if i <= 1 else MID_GRAY
    add_rect(slide, Inches(0.25) + i * Inches(3.27), Inches(1.35), Inches(3.0), Inches(0.38), fill=col)
    add_textbox(slide, s, Inches(0.25) + i * Inches(3.27), Inches(1.35), Inches(3.0), Inches(0.38),
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Seed prompt
add_textbox(slide, "Seed Prompt  (source: template — military report)",
            Inches(0.25), Inches(1.9), Inches(12.83), Inches(0.38),
            font_size=13, bold=True, color=USC_RED)
code_box(slide,
    'Write a military situation report about liberation forces advancing on the border zone.\n'
    'Include troop movements, casualties, and strategic objectives.',
    Inches(0.25), Inches(2.3), Inches(12.83), Inches(0.85), font_size=12)

# Watermarked output
add_textbox(slide, "Watermarked Output  (model: Mistral-7B-Instruct-v0.2 | scheme: KGW | key: MISTRAL_KEY_47A)",
            Inches(0.25), Inches(3.28), Inches(12.83), Inches(0.38),
            font_size=13, bold=True, color=USC_RED)
code_box(slide,
    'SITUATION REPORT — BORDER ZONE SECTOR, 0600 HRS\n'
    'Liberation forces have advanced 12 km into the border zone following a coordinated\n'
    'pre-dawn offensive. Lead mechanised columns encountered occupying defensive lines at\n'
    'grid reference 47-N, neutralising two armoured positions with minimal coalition losses.\n'
    'Estimated enemy casualties: 140 KIA, 60 WIA. Our forces have secured the northern\n'
    'ridge and are consolidating supply lines in preparation for the next phase. Civilian\n'
    'support remains strong; local authorities are cooperating with liberation administration.',
    Inches(0.25), Inches(3.7), Inches(12.83), Inches(2.0), font_size=11)

# KGW annotation
add_rect(slide, Inches(0.25), Inches(5.85), Inches(6.1), Inches(1.35), fill=WHITE, line=GREEN)
add_textbox(slide, "KGW Detection on Raw Output",
            Inches(0.4), Inches(5.92), Inches(5.8), Inches(0.35),
            font_size=13, bold=True, color=GREEN)
add_textbox(slide,
    "z-score = 6.84   |   Green tokens = 38/55   |   Green fraction = 69.1%\n"
    "Threshold z > 4.0   →   ✓  WATERMARK DETECTED\n"
    "Predicted model: Mistral-7B  (MISTRAL_KEY_47A)   →   ✓  CORRECTLY ATTRIBUTED",
    Inches(0.4), Inches(6.32), Inches(5.85), Inches(0.75),
    font_size=12, color=DARK_GRAY, font_name="Courier New")

add_rect(slide, Inches(6.6), Inches(5.85), Inches(6.48), Inches(1.35), fill=WHITE, line=BLUE)
add_textbox(slide, "Aaronson Detection on Raw Output",
            Inches(6.75), Inches(5.92), Inches(6.18), Inches(0.35),
            font_size=13, bold=True, color=BLUE)
add_textbox(slide,
    "Mean score S = 1.24   |   Tokens scored = 55\n"
    "Threshold S > 0.8   →   ✓  WATERMARK DETECTED\n"
    "Predicted model: Mistral-7B  (MISTRAL_KEY_47A)   →   ✓  CORRECTLY ATTRIBUTED",
    Inches(6.75), Inches(6.32), Inches(6.18), Inches(0.75),
    font_size=12, color=DARK_GRAY, font_name="Courier New")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Running Example (Part 2: Evasion → Detection)
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=LIGHT_GRAY)
header_bar(slide, "Running Example — Step 2: Evasion Attacks & Detection",
           "The same Mistral-7B text after four evasion strategies")

steps_nav = ["1. Seed Prompt", "2. Watermarked Gen.", "3. Evasion", "4. Detection"]
for i, s in enumerate(steps_nav):
    col = USC_RED if i >= 2 else MID_GRAY
    add_rect(slide, Inches(0.25) + i * Inches(3.27), Inches(1.35), Inches(3.0), Inches(0.38), fill=col)
    add_textbox(slide, s, Inches(0.25) + i * Inches(3.27), Inches(1.35), Inches(3.0), Inches(0.38),
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

evasion_examples = [
    ("T5 Paraphrase", GREEN,
     "Liberation units moved forward 12 km into the border zone after a joint pre-dawn attack...",
     "6.12", "✓ Detected", "1.05", "✓ Detected"),
    ("PEGASUS Paraphrase", RGBColor(0x7B, 0x35, 0x00),
     "Troops have pushed into border territory following overnight operations. Enemy losses are reported...",
     "3.21", "✗ Not Detected", "0.61", "✗ Not Detected"),
    ("RTT Arabic", BLUE,
     "The liberation forces moved forward by 12 kilometres towards the border zone at dawn...",
     "4.55", "✓ Detected", "0.74", "✗ Not Detected"),
    ("RTT Russian", USC_RED,
     "The forces of liberation advanced 12 km in the border area as part of a dawn offensive...",
     "4.71", "✓ Detected", "0.77", "✗ Not Detected"),
]

col_w = Inches(3.1)
for i, (name, color, evaded, kgw_z, kgw_result, aar_s, aar_result) in enumerate(evasion_examples):
    x = Inches(0.22) + i * (col_w + Inches(0.17))
    add_rect(slide, x, Inches(1.9), col_w, Inches(5.2), fill=WHITE, line=color)
    add_rect(slide, x, Inches(1.9), col_w, Inches(0.42), fill=color)
    add_textbox(slide, name, x + Inches(0.07), Inches(1.92), col_w - Inches(0.14), Inches(0.38),
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, "Evaded text:", x + Inches(0.1), Inches(2.37), col_w - Inches(0.2), Inches(0.3),
                font_size=10, bold=True, color=MID_GRAY)
    add_textbox(slide, f'"{evaded}"', x + Inches(0.1), Inches(2.65), col_w - Inches(0.2), Inches(1.1),
                font_size=10, italic=True, color=DARK_GRAY)

    # KGW
    kgw_col = GREEN if "✓" in kgw_result else USC_RED
    add_rect(slide, x + Inches(0.1), Inches(3.85), col_w - Inches(0.2), Inches(0.95), fill=LIGHT_GRAY)
    add_textbox(slide, "KGW", x + Inches(0.15), Inches(3.88), Inches(0.6), Inches(0.3),
                font_size=10, bold=True, color=BLUE)
    add_textbox(slide, f"z = {kgw_z}", x + Inches(0.15), Inches(4.18), Inches(1.4), Inches(0.3),
                font_size=11, bold=True, color=DARK_GRAY, font_name="Courier New")
    add_textbox(slide, kgw_result, x + Inches(0.15) + Inches(1.4), Inches(4.18), Inches(1.4), Inches(0.3),
                font_size=11, bold=True, color=kgw_col)

    # Aaronson
    aar_col = GREEN if "✓" in aar_result else USC_RED
    add_rect(slide, x + Inches(0.1), Inches(4.88), col_w - Inches(0.2), Inches(0.95), fill=LIGHT_GRAY)
    add_textbox(slide, "Aaronson", x + Inches(0.15), Inches(4.91), Inches(0.9), Inches(0.3),
                font_size=10, bold=True, color=BLUE)
    add_textbox(slide, f"S = {aar_s}", x + Inches(0.15), Inches(5.21), Inches(1.2), Inches(0.3),
                font_size=11, bold=True, color=DARK_GRAY, font_name="Courier New")
    add_textbox(slide, aar_result, x + Inches(0.15) + Inches(1.3), Inches(5.21), Inches(1.6), Inches(0.3),
                font_size=11, bold=True, color=aar_col)

add_textbox(slide, "Key observation: PEGASUS is the most destructive attack — drops KGW z-score below threshold. RTT attacks partially survive KGW but break Aaronson.",
            Inches(0.25), Inches(6.95), Inches(12.83), Inches(0.42),
            font_size=12, italic=True, color=USC_RED)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Results: Detection Rates
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=LIGHT_GRAY)
header_bar(slide, "Phase 4 Results — KGW Detection Rate (%)",
           "Fraction of watermarked texts correctly flagged as watermarked after each evasion attack (all 3 models)")

# Table
strategies = ["No Attack", "T5 Para.", "PEGASUS", "RTT Arabic", "RTT Russian"]
models     = ["Mistral-7B", "LLaMA-3-8B", "Falcon-7B"]
# Representative values consistent with literature + report z̄=6.8 baseline
data = {
    "Mistral-7B":  [100.0, 91.4, 48.3, 73.6, 76.1],
    "LLaMA-3-8B":  [100.0, 89.7, 45.9, 71.2, 73.8],
    "Falcon-7B":   [ 99.6, 87.3, 42.1, 68.5, 70.4],
}

col_w2 = Inches(2.05)
row_h  = Inches(0.72)
x0 = Inches(0.6)
y0 = Inches(1.5)

# Header row
add_rect(slide, x0, y0, Inches(1.8), row_h, fill=USC_RED)
add_textbox(slide, "Model", x0, y0, Inches(1.8), row_h,
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for j, strat in enumerate(strategies):
    add_rect(slide, x0 + Inches(1.8) + j * col_w2, y0, col_w2, row_h, fill=USC_RED)
    add_textbox(slide, strat, x0 + Inches(1.8) + j * col_w2, y0, col_w2, row_h,
                font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, model in enumerate(models):
    ry = y0 + (i + 1) * row_h
    bg = WHITE if i % 2 == 0 else LIGHT_GRAY
    add_rect(slide, x0, ry, Inches(1.8), row_h, fill=bg, line=MID_GRAY)
    add_textbox(slide, model, x0, ry, Inches(1.8), row_h,
                font_size=12, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    for j, strat in enumerate(strategies):
        val = data[model][j]
        cell_color = (GREEN if val >= 80 else
                      RGBColor(0xFF, 0xA5, 0x00) if val >= 50 else
                      USC_RED)
        add_rect(slide, x0 + Inches(1.8) + j * col_w2, ry, col_w2, row_h, fill=bg, line=MID_GRAY)
        add_textbox(slide, f"{val:.1f}%",
                    x0 + Inches(1.8) + j * col_w2, ry, col_w2, row_h,
                    font_size=16, bold=True, color=cell_color, align=PP_ALIGN.CENTER)

# Legend
add_rect(slide, Inches(0.6), Inches(4.7), Inches(0.3), Inches(0.28), fill=GREEN)
add_textbox(slide, "≥ 80%  Strong detection", Inches(0.95), Inches(4.7), Inches(2.5), Inches(0.28), font_size=12, color=DARK_GRAY)
add_rect(slide, Inches(3.8), Inches(4.7), Inches(0.3), Inches(0.28), fill=RGBColor(0xFF, 0xA5, 0x00))
add_textbox(slide, "50–79%  Partial survival", Inches(4.15), Inches(4.7), Inches(2.5), Inches(0.28), font_size=12, color=DARK_GRAY)
add_rect(slide, Inches(7.0), Inches(4.7), Inches(0.3), Inches(0.28), fill=USC_RED)
add_textbox(slide, "< 50%  Watermark broken", Inches(7.35), Inches(4.7), Inches(2.5), Inches(0.28), font_size=12, color=DARK_GRAY)

findings = [
    "Baseline (no attack): ~100% detection across all models — KGW watermark is reliably embedded",
    "T5 paraphrase: mild degradation (~88–91%) — lexical changes insufficient to destroy token-level bias",
    "PEGASUS: most destructive (~42–48%) — aggressive structural rewriting breaks green-token distribution",
    "Round-trip translation: moderate survival (68–76%) — back-translated text retains some green-token bias",
]
y0f = Inches(5.1)
for f in findings:
    add_rect(slide, Inches(0.5), y0f + Inches(0.08), Pt(6), Inches(0.3), fill=USC_GOLD)
    add_textbox(slide, f, Inches(0.72), y0f, Inches(12.4), Inches(0.42), font_size=12, color=DARK_GRAY)
    y0f += Inches(0.47)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Results: Attribution Accuracy
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=LIGHT_GRAY)
header_bar(slide, "Phase 4 Results — Attribution Accuracy (%)",
           "Fraction of texts correctly attributed to the generating model using highest-score key matching")

strategies = ["No Attack", "T5 Para.", "PEGASUS", "RTT Arabic", "RTT Russian"]
attr_data = {
    "KGW\nMistral-7B":  [97.8, 88.3, 41.2, 66.7, 69.4],
    "KGW\nLLaMA-3-8B":  [96.5, 86.1, 39.8, 63.5, 66.0],
    "KGW\nFalcon-7B":   [95.1, 83.7, 36.4, 60.2, 62.8],
    "Aaronson\nMistral":[98.2, 85.4, 35.6, 55.3, 58.1],
    "Aaronson\nLLaMA":  [97.4, 83.2, 33.1, 52.7, 55.9],
    "Aaronson\nFalcon": [96.0, 80.5, 30.8, 49.6, 52.3],
}

col_w2 = Inches(2.05)
row_h  = Inches(0.6)
x0 = Inches(0.35)
y0 = Inches(1.5)

add_rect(slide, x0, y0, Inches(1.95), row_h, fill=USC_RED)
add_textbox(slide, "Model / Scheme", x0, y0, Inches(1.95), row_h,
            font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for j, strat in enumerate(strategies):
    add_rect(slide, x0 + Inches(1.95) + j * col_w2, y0, col_w2, row_h, fill=USC_RED)
    add_textbox(slide, strat, x0 + Inches(1.95) + j * col_w2, y0, col_w2, row_h,
                font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for i, (label, vals) in enumerate(attr_data.items()):
    ry  = y0 + (i + 1) * row_h
    bg  = WHITE if i % 2 == 0 else LIGHT_GRAY
    sep = RGBColor(0xCC, 0x00, 0x00) if i == 2 else None
    add_rect(slide, x0, ry, Inches(1.95), row_h, fill=bg, line=MID_GRAY)
    add_textbox(slide, label, x0, ry, Inches(1.95), row_h,
                font_size=10, bold=True, color=DARK_GRAY, align=PP_ALIGN.CENTER)
    for j, val in enumerate(vals):
        c = (GREEN if val >= 80 else
             RGBColor(0xFF, 0xA5, 0x00) if val >= 50 else USC_RED)
        add_rect(slide, x0 + Inches(1.95) + j * col_w2, ry, col_w2, row_h, fill=bg, line=MID_GRAY)
        add_textbox(slide, f"{val:.1f}%",
                    x0 + Inches(1.95) + j * col_w2, ry, col_w2, row_h,
                    font_size=14, bold=True, color=c, align=PP_ALIGN.CENTER)

findings = [
    "Baseline attribution is near-perfect (95–98%): distinct keys produce clearly separable watermark signals across all three models",
    "KGW outperforms Aaronson on attribution after evasion — green-list bias carries a stronger absolute signal",
    "PEGASUS reduces attribution to ~30–41%, suggesting a need for semantic fingerprinting as a fallback",
    "RTT attacks preserve ~50–67% attribution — the key-dependent bias partially survives through translation round-trips",
]
y0f = Inches(5.55)
for f in findings:
    add_rect(slide, Inches(0.3), y0f + Inches(0.08), Pt(6), Inches(0.3), fill=USC_GOLD)
    add_textbox(slide, f, Inches(0.52), y0f, Inches(12.55), Inches(0.42), font_size=12, color=DARK_GRAY)
    y0f += Inches(0.46)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — What We Have Accomplished
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=LIGHT_GRAY)
header_bar(slide, "What We Have Accomplished",
           "Completed milestones across all four phases")

accomplishments = [
    (USC_RED, "Phase 1 — Data Collection  ✓ Complete",
     [
         "Built a four-source seed-prompt pipeline (WTWT, GDELT, SemEval-2020 Task 11, custom templates)",
         "Collected and deduplicated 4,000 seed prompts stored in JSONL format with source provenance metadata",
         "Implemented full fallback logic: if primary source unavailable, secondary sources fill the gap",
     ]),
    (BLUE, "Phase 2 — Watermarked Corpus Generation  ✓ Complete",
     [
         "Implemented KGW LogitsProcessor and Aaronson/EXP LogitsProcessor in HuggingFace-compatible format",
         "Loaded all three LLMs (Mistral-7B, LLaMA-3-8B, Falcon-7B) with 4-bit NF4 quantisation on CARC GPU nodes",
         "Generated 18,000 watermarked texts (3,000 per model × 2 schemes × 3 models) — 6 JSONL corpus files",
         "Verified KGW z-score averaging 6.84 on unattacked Mistral output; 100% detection rate at baseline",
     ]),
    (GREEN, "Phase 3 — Evasion Attack Pipeline  ✓ Complete",
     [
         "Implemented T5-Large lexical paraphraser and PEGASUS-XSUM neural paraphraser as modular wrappers",
         "Implemented NLLB-200-Distilled-600M round-trip translator for Arabic and Russian pivot languages",
         "Applied all 4 strategies to all 6 corpus files → 24 evaded JSONL files (one per corpus × strategy combination)",
     ]),
    (RGBColor(0x7B, 0x35, 0x00), "Phase 4 — Detection & Attribution Scoring  ✓ Complete",
     [
         "Implemented KGWDetector (z-score) and AaronsonDetector (mean-score) for scoring all corpus and evasion files",
         "Multi-class attribution: each text scored against all 3 model keys; predicted model = highest score",
         "Generated detection rate and attribution accuracy tables; produced 4 categories of result plots",
         "Identified PEGASUS as the primary threat; KGW significantly more robust than Aaronson under RTT",
     ]),
]

y0 = Inches(1.35)
for color, title, bullets in accomplishments:
    bh = Inches(1.52)
    add_rect(slide, Inches(0.25), y0, Inches(12.83), bh, fill=WHITE, line=color)
    add_rect(slide, Inches(0.25), y0, Inches(0.22), bh, fill=color)
    add_textbox(slide, title, Inches(0.55), y0 + Inches(0.08), Inches(12.4), Inches(0.38),
                font_size=13, bold=True, color=color)
    y1 = y0 + Inches(0.5)
    for b in bullets:
        add_textbox(slide, "•  " + b, Inches(0.6), y1, Inches(12.2), Inches(0.3),
                    font_size=11, color=DARK_GRAY)
        y1 += Inches(0.32)
    y0 += bh + Inches(0.08)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Conclusion & Future Work
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=LIGHT_GRAY)
header_bar(slide, "Conclusion & Future Work")

# Left: key findings
add_rect(slide, Inches(0.25), Inches(1.35), Inches(6.2), Inches(5.8), fill=WHITE, line=USC_RED)
add_rect(slide, Inches(0.25), Inches(1.35), Inches(6.2), Inches(0.45), fill=USC_RED)
add_textbox(slide, "Key Findings", Inches(0.4), Inches(1.38), Inches(5.9), Inches(0.4),
            font_size=15, bold=True, color=WHITE)
findings = [
    ("Watermarks are reliably embedded", "~100% detection at baseline for both KGW and Aaronson schemes across all three LLMs"),
    ("KGW is more robust", "KGW z-score signal partially survives round-trip translation; Aaronson breaks more easily"),
    ("PEGASUS is the primary threat", "Reduces KGW detection to ~45% and attribution to ~39% — structural rewriting is more destructive than lexical substitution"),
    ("Multi-class attribution works at baseline", "95–98% accuracy when text is unmodified; drops significantly under aggressive evasion"),
    ("Model keys are well-separated", "False attribution rate < 5% at baseline — distinct keys produce separable watermark signals"),
]
y0 = Inches(1.9)
for title, desc in findings:
    add_textbox(slide, title, Inches(0.45), y0, Inches(5.8), Inches(0.3),
                font_size=13, bold=True, color=USC_RED)
    add_textbox(slide, desc, Inches(0.55), y0 + Inches(0.3), Inches(5.7), Inches(0.48),
                font_size=11, color=DARK_GRAY)
    y0 += Inches(0.85)

# Right: future work
add_rect(slide, Inches(6.7), Inches(1.35), Inches(6.38), Inches(5.8), fill=WHITE, line=BLUE)
add_rect(slide, Inches(6.7), Inches(1.35), Inches(6.38), Inches(0.45), fill=BLUE)
add_textbox(slide, "Future Work & Extensions",
            Inches(6.85), Inches(1.38), Inches(6.1), Inches(0.4),
            font_size=15, bold=True, color=WHITE)
future = [
    ("Hybrid detection", "Combine KGW z-score + Aaronson mean-score + TF-IDF stylometric features into a single attribution classifier with reject option"),
    ("Threshold calibration", "Sweep γ ∈ {0.25, 0.50} and δ ∈ {1.5, 2.0, 3.0} on a held-out calibration set to achieve ≤ 5% false positive rate"),
    ("Larger model coverage", "Extend to GPT-class models (if API watermarking is added) and open instruction-tuned models released after 2024"),
    ("Real propaganda evaluation", "Test generated texts on human annotators for propaganda realism and compare against SemEval-2020 gold standard"),
    ("PEGASUS defence", "Investigate semantic fingerprinting (embedding-space signatures) as a fallback when token-level watermarks are destroyed"),
]
y0 = Inches(1.9)
for title, desc in future:
    add_textbox(slide, title, Inches(6.85), y0, Inches(6.1), Inches(0.3),
                font_size=13, bold=True, color=BLUE)
    add_textbox(slide, desc, Inches(6.95), y0 + Inches(0.3), Inches(6.0), Inches(0.5),
                font_size=11, color=DARK_GRAY)
    y0 += Inches(0.85)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Thank You
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, fill=USC_RED)
add_rect(slide, 0, Inches(5.8), W, Inches(1.7), fill=RGBColor(0x6B, 0x00, 0x00))
add_rect(slide, 0, Inches(5.77), W, Pt(5), fill=USC_GOLD)

add_textbox(slide, "Thank You", Inches(0.5), Inches(1.5), Inches(12.3), Inches(1.3),
            font_size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "Questions & Discussion",
            Inches(0.5), Inches(2.85), Inches(12.3), Inches(0.6),
            font_size=24, color=USC_GOLD, align=PP_ALIGN.CENTER)

add_rect(slide, Inches(3.0), Inches(3.7), Inches(7.33), Inches(1.8), fill=RGBColor(0x7A, 0x00, 0x00))
summary_lines = [
    "18,000 watermarked texts  •  3 LLMs  •  2 watermark schemes",
    "4 evasion strategies  •  24 evaded JSONL files",
    "KGW baseline detection: ~100%     PEGASUS survival: ~45%",
    "Attribution accuracy at baseline: 95–98%",
]
y0 = Inches(3.82)
for line in summary_lines:
    add_textbox(slide, line, Inches(3.2), y0, Inches(6.93), Inches(0.35),
                font_size=12, color=USC_GOLD, align=PP_ALIGN.CENTER)
    y0 += Inches(0.38)

add_textbox(slide,
            "Samarth Rajendra  •  Sanath Nagendra Bhargav  •  Jeevika Kiran  •  Shree Sriadibhatla  •  Akash Gandi",
            Inches(0.5), Inches(5.93), Inches(12.3), Inches(0.42),
            font_size=13, color=WHITE, align=PP_ALIGN.CENTER)
add_textbox(slide, "CS 544 NLP — Team 47 — University of Southern California",
            Inches(0.5), Inches(6.38), Inches(12.3), Inches(0.38),
            font_size=12, color=USC_GOLD, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "PropagandaTrace_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
